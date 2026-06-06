#!/usr/bin/env python3
# ==============================================================================
# BLM3810 Biyoenformatiğe Giriş - Proje Sunumu Üreteci
# Stil: 333F70 lacivert, Open Sans / Unbounded, 16:9
# Çıktı: presentation/Biyoenformatik_Sunum.pptx
# ==============================================================================

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

BASE = os.getcwd()
FIG = os.path.join(BASE, "results", "figures")
OUT = os.path.join(BASE, "presentation", "Biyoenformatik_Sunum.pptx")

# --- Renk paleti ---
NAVY   = RGBColor(0x33, 0x3F, 0x70)   # ana marka rengi
NAVY_D = RGBColor(0x24, 0x2C, 0x50)   # koyu lacivert
RED    = RGBColor(0xE6, 0x4B, 0x35)   # tümör
GREEN  = RGBColor(0x00, 0xA0, 0x87)   # non-tümör
BLUE   = RGBColor(0x44, 0x72, 0xC4)   # accent
LBLUE  = RGBColor(0x5B, 0x9B, 0xD5)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREY   = RGBColor(0x59, 0x59, 0x59)
LGREY  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x22, 0x22, 0x22)

F_TITLE = "Unbounded"     # büyük başlıklar (örnek sunumdaki gibi)
F_BODY  = "Open Sans"     # gövde

# --- Slayt boyutu 16:9 (16x9 inch) ---
prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)
BLANK = prs.slide_layouts[6]

SW, SH = Inches(16), Inches(9)


# ==============================================================================
# Yardımcı fonksiyonlar
# ==============================================================================

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, x, y, w, h, fill, line=None, line_w=None, shadow=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp

def rrect(slide, x, y, w, h, fill, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w or Pt(1.25)
    sp.shadow.inherit = False
    return sp

def txt(slide, x, y, w, h, text, size, color=DARK, bold=False, font=F_BODY,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_sp=1.0, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    if line_sp: p.line_spacing = line_sp
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb

def bullets(slide, x, y, w, h, items, size=20, color=DARK, gap=8,
            bullet_color=NAVY, font=F_BODY, line_sp=1.05):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            text, lvl = it
        else:
            text, lvl = it, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = line_sp
        p.space_after = Pt(gap)
        p.level = lvl
        mark = "•  " if lvl == 0 else "–  "
        r1 = p.add_run(); r1.text = mark
        r1.font.size = Pt(size); r1.font.bold = True
        r1.font.color.rgb = bullet_color; r1.font.name = font
        r2 = p.add_run(); r2.text = text
        r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = font
    return tb

def content_header(slide, title, kicker=None):
    """Beyaz içerik slaytı başlığı: sol accent + başlık + alt çizgi."""
    rect(slide, Inches(0), Inches(0), SW, Inches(1.55), WHITE)
    rect(slide, Inches(0), Inches(0), Inches(0.28), Inches(1.55), NAVY)
    if kicker:
        txt(slide, Inches(0.7), Inches(0.25), Inches(13), Inches(0.4),
            kicker, 14, NAVY, bold=True, font=F_BODY)
        ty = Inches(0.6)
    else:
        ty = Inches(0.42)
    txt(slide, Inches(0.7), ty, Inches(14.6), Inches(0.85),
        title, 30, NAVY, bold=True, font=F_TITLE)
    rect(slide, Inches(0.72), Inches(1.42), Inches(2.2), Pt(3), RED)

def footer(slide, n):
    txt(slide, Inches(0.5), Inches(8.55), Inches(10), Inches(0.35),
        "GSE76427 · Hepatosellüler Karsinom · BLM3810", 11, GREY, font=F_BODY)
    txt(slide, Inches(14.8), Inches(8.55), Inches(0.8), Inches(0.35),
        str(n), 12, GREY, bold=True, align=PP_ALIGN.RIGHT, font=F_BODY)

def pic_fit(slide, path, x, y, max_w, max_h, center_x=True):
    """Resmi oranını koruyarak max kutuya sığdır."""
    img = Image.open(path); iw, ih = img.size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w; h = Emu(int(max_w / ar))
    else:
        h = max_h; w = Emu(int(max_h * ar))
    px = x + Emu(int((max_w - w) / 2)) if center_x else x
    py = y + Emu(int((max_h - h) / 2))
    slide.shapes.add_picture(path, px, py, width=w, height=h)
    return w, h

def card(slide, x, y, w, h, header, body, hcolor, body_size=15, head_size=16):
    """Renkli başlıklı kart."""
    rrect(slide, x, y, w, h, LGREY)
    rect(slide, x, y, w, Inches(0.62), hcolor)
    txt(slide, x+Inches(0.12), y+Inches(0.06), w-Inches(0.24), Inches(0.5),
        header, head_size, WHITE, bold=True, font=F_BODY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, x+Inches(0.2), y+Inches(0.74), w-Inches(0.4), h-Inches(0.84),
        body, body_size, DARK, font=F_BODY, line_sp=1.05)


# ==============================================================================
# SLAYT 1 — KAPAK
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, NAVY)
# sağ accent şerit
rect(s, Inches(13.6), Inches(0), Inches(2.4), SH, NAVY_D)
rect(s, Inches(13.55), Inches(0), Inches(0.08), SH, RED)
# kicker
txt(s, Inches(1.1), Inches(1.6), Inches(12), Inches(0.5),
    "BLM3810 BİYOENFORMATİĞE GİRİŞ — DÖNEM PROJESİ", 16, LBLUE, bold=True, font=F_BODY)
# başlık
txt(s, Inches(1.0), Inches(2.4), Inches(12.2), Inches(2.6),
    "Gen İfadesi Analizleri ve Makine Öğrenmesi ile Hepatosellüler Karsinom Sınıflandırması",
    40, WHITE, bold=True, font=F_TITLE, line_sp=1.05)
# alt çizgi
rect(s, Inches(1.05), Inches(5.15), Inches(3.2), Pt(4), RED)
# künye
info = [
    ("Hazırlayan:", "Mert GÜLER"),
    ("Kurum:", "Yıldız Teknik Üniversitesi – Bilgisayar Mühendisliği Bölümü"),
    ("Ders:", "BLM3810 Biyoenformatiğe Giriş"),
    ("Öğretim Üyesi:", "Dr. S. Sevgi TURGUT ÖGME"),
    ("Veri Seti:", "GSE76427 (NCBI GEO)"),
    ("Tarih:", "Haziran 2026"),
]
yy = Inches(5.55)
for k, v in info:
    txt(s, Inches(1.1), yy, Inches(2.3), Inches(0.4), k, 15, LBLUE, bold=True, font=F_BODY)
    txt(s, Inches(3.4), yy, Inches(9.5), Inches(0.4), v, 15, WHITE, font=F_BODY)
    yy += Inches(0.46)


# ==============================================================================
# SLAYT 2 — SUNUM PLANI
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "Sunum Planı")
plan = [
    ("Giriş: Biyoenformatik, gen ifadesi ve HCC", BLUE),
    ("Veri Seti ve Ön İşleme", GREEN),
    ("Diferansiyel Gen İfadesi (DGE) Analizi", RED),
    ("Sağkalım Analizi ve Cox Modeli", ORANGE),
    ("WGCNA Ağ Analizi ve Hub Genler", LBLUE),
    ("GO Zenginleştirme Analizi", NAVY),
    ("Makine Öğrenmesi ile Sınıflandırma", BLUE),
    ("Sonuçlar ve Değerlendirme", GREEN),
]
yy = Inches(2.1)
for i, (t, c) in enumerate(plan, 1):
    rrect(s, Inches(1.2), yy, Inches(0.62), Inches(0.62), c)
    txt(s, Inches(1.2), yy, Inches(0.62), Inches(0.62), str(i), 20, WHITE,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_TITLE)
    txt(s, Inches(2.1), yy, Inches(12), Inches(0.62), t, 21, DARK,
        anchor=MSO_ANCHOR.MIDDLE, font=F_BODY)
    yy += Inches(0.78)
footer(s, 2)


# ==============================================================================
# SLAYT 3 — GİRİŞ
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "Giriş: Biyoenformatik ve HCC", kicker="MOTİVASYON")
bullets(s, Inches(0.8), Inches(2.0), Inches(8.2), Inches(6),
    [
     "Gen ifadesi: bir genin ürettiği mRNA miktarı; hücrenin hangi genleri ne kadar aktif kullandığını gösterir.",
     "Kanser, gen ifadesi profillerindeki belirgin değişimlerle karakterizedir.",
     ("Onkogenler aşırı aktive olur (up-regulated)", 1),
     ("Tümör süpresörler baskılanır (down-regulated)", 1),
     "Hepatosellüler Karsinom (HCC): en yaygın primer karaciğer kanseri.",
     "Karaciğer; detoksifikasyon, metabolizma ve protein sentezinin merkezidir — tümörde bu fonksiyonlar bozulur.",
     "Amaç: Tümör/normal dokuyu ayıran genleri bulmak ve makine öğrenmesi ile sınıflandırmak.",
    ], size=19, gap=10)
# sağ kutu — proje akışı
rrect(s, Inches(9.4), Inches(2.0), Inches(5.9), Inches(5.6), LGREY)
txt(s, Inches(9.6), Inches(2.15), Inches(5.5), Inches(0.5),
    "İki Aşamalı Proje", 18, NAVY, bold=True, font=F_BODY)
flow = ["1 · Veri Ön İşleme", "2 · DGE Analizi", "3 · Sağkalım & Cox",
        "4 · WGCNA Ağ Analizi", "5 · GO Zenginleştirme", "6 · ML Sınıflandırma"]
yy = Inches(2.75)
cols = [GREEN, RED, ORANGE, LBLUE, NAVY, BLUE]
for f, c in zip(flow, cols):
    rrect(s, Inches(9.7), yy, Inches(5.3), Inches(0.62), c)
    txt(s, Inches(9.7), yy, Inches(5.3), Inches(0.62), f, 16, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_BODY)
    yy += Inches(0.78)
footer(s, 3)


# ==============================================================================
# SLAYT 4 — VERİ SETİ
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "Veri Seti: GSE76427", kicker="NCBI GEO")
# sol bilgi kartları
stats = [
    ("167", "Toplam Örnek", NAVY),
    ("115", "Tümör Dokusu", RED),
    ("52", "Non-Tümör Dokusu", GREEN),
    ("47.322", "Ham Prob", BLUE),
]
xx = Inches(0.8)
for val, lab, c in stats:
    rrect(s, xx, Inches(2.0), Inches(3.35), Inches(1.6), c)
    txt(s, xx, Inches(2.15), Inches(3.35), Inches(0.9), val, 40, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_TITLE)
    txt(s, xx, Inches(3.05), Inches(3.35), Inches(0.5), lab, 15, WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_BODY)
    xx += Inches(3.55)
# detay tablo
det = [
    ("Hastalık", "Hepatosellüler Karsinom (Karaciğer Kanseri)"),
    ("Platform", "Illumina HumanHT-12 V4.0 BeadChip (GPL10558)"),
    ("Sınıflandırma Etiketi", "Tümör vs Non-Tümör (kanser evresi DEĞİL)"),
    ("Sağkalım Verisi", "OS: 115 örnek (23 olay) · RFS: 108 örnek (48 olay)"),
    ("Klinik Değişkenler", "Yaş (14–93), cinsiyet (93E/22K), BCLC/TNM evre"),
]
yy = Inches(4.0)
for i, (k, v) in enumerate(det):
    bg = LGREY if i % 2 == 0 else WHITE
    rect(s, Inches(0.8), yy, Inches(14.4), Inches(0.62), bg)
    txt(s, Inches(1.0), yy, Inches(4.2), Inches(0.62), k, 16, NAVY, bold=True,
        anchor=MSO_ANCHOR.MIDDLE, font=F_BODY)
    txt(s, Inches(5.3), yy, Inches(9.7), Inches(0.62), v, 16, DARK,
        anchor=MSO_ANCHOR.MIDDLE, font=F_BODY)
    yy += Inches(0.62)
txt(s, Inches(0.8), yy+Inches(0.15), Inches(14), Inches(0.5),
    "Seçim gerekçesi: 2 karşılaştırılabilir sınıf + survival metadata + evre etiketi kullanılmaması → proje kurallarına uygun.",
    15, GREY, italic=True, font=F_BODY)
footer(s, 4)


# ==============================================================================
# SLAYT 5 — ÖN İŞLEME
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "Veri Ön İşleme Hattı", kicker="PIPELINE")
steps = [
    ("Ham Veri", "47.322 prob\nNon-normalized\nAVG_Signal", BLUE),
    ("Normalizasyon", "log2 dönüşümü\n+ Quantile norm\n(limma)", GREEN),
    ("Prob Filtre", "Detection p<0.05\n+ varyans filtre\n→ 24.345 prob", ORANGE),
    ("Anotasyon", "Prob → gen sembolü\nDuplicate: max var.\n→ 14.323 gen", LBLUE),
    ("HVG Seçimi", "Top-3000\nyüksek varyanslı\ngen", RED),
]
xx = Inches(0.7)
cw = Inches(2.75)
for i, (h, b, c) in enumerate(steps):
    card(s, xx, Inches(2.0), cw, Inches(2.4), h, b, c, body_size=14)
    if i < len(steps)-1:
        txt(s, xx+cw-Inches(0.05), Inches(2.85), Inches(0.4), Inches(0.6),
            "→", 28, NAVY, bold=True, align=PP_ALIGN.CENTER, font=F_BODY)
    xx += cw + Inches(0.18)
# normalizasyon öncesi/sonrası görsel
txt(s, Inches(0.8), Inches(4.7), Inches(14), Inches(0.4),
    "Normalizasyon Öncesi / Sonrası — örnek dağılımları hizalanır:", 16, NAVY, bold=True, font=F_BODY)
pic_fit(s, os.path.join(FIG, "boxplot_before_norm.png"), Inches(0.8), Inches(5.15), Inches(7.0), Inches(3.0))
pic_fit(s, os.path.join(FIG, "boxplot_after_norm.png"), Inches(8.1), Inches(5.15), Inches(7.0), Inches(3.0))
footer(s, 5)


# ==============================================================================
# SLAYT 6 — KEŞİFSEL ANALİZ
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "Keşifsel Veri Analizi: PCA & Isı Haritası", kicker="EDA")
pic_fit(s, os.path.join(FIG, "pca_colored_by_group.png"), Inches(0.7), Inches(2.0), Inches(7.3), Inches(5.4))
pic_fit(s, os.path.join(FIG, "heatmap_top20_hvg.png"), Inches(8.2), Inches(2.0), Inches(7.1), Inches(5.4))
txt(s, Inches(0.7), Inches(7.5), Inches(14.5), Inches(0.9),
    "Tümör (kırmızı) ve non-tümör (yeşil) örnekler PCA ve ısı haritasında belirgin biçimde ayrı kümeler oluşturur → "
    "gen ifadesi düzeyinde güçlü ayırt edici sinyal mevcuttur.",
    16, GREY, italic=True, font=F_BODY, align=PP_ALIGN.CENTER)
footer(s, 6)


# ==============================================================================
# SLAYT 7 — DGE YÖNTEM + SONUÇ
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "Diferansiyel Gen İfadesi (DGE)", kicker="limma")
# sol: yöntem + sayılar
bullets(s, Inches(0.8), Inches(1.95), Inches(6.6), Inches(2.4),
    [
     "limma, tüm 14.323 gen üzerinde",
     "Paired tasarım: hasta bloğu (duplicateCorrelation)",
     "Kontrast: Tumor − NonTumor",
     "Eşik: adj.P.Val < 0.05  ve  |logFC| > 1",
    ], size=17, gap=7)
# sayı kutuları
rrect(s, Inches(0.8), Inches(4.0), Inches(2.05), Inches(1.3), NAVY)
txt(s, Inches(0.8), Inches(4.1), Inches(2.05), Inches(0.8), "705", 34, WHITE, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_TITLE)
txt(s, Inches(0.8), Inches(4.85), Inches(2.05), Inches(0.4), "Toplam DEG", 13, WHITE,
    align=PP_ALIGN.CENTER, font=F_BODY)
rrect(s, Inches(3.0), Inches(4.0), Inches(2.05), Inches(1.3), RED)
txt(s, Inches(3.0), Inches(4.1), Inches(2.05), Inches(0.8), "216", 34, WHITE, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_TITLE)
txt(s, Inches(3.0), Inches(4.85), Inches(2.05), Inches(0.4), "Up ↑", 13, WHITE,
    align=PP_ALIGN.CENTER, font=F_BODY)
rrect(s, Inches(5.2), Inches(4.0), Inches(2.05), Inches(1.3), GREEN)
txt(s, Inches(5.2), Inches(4.1), Inches(2.05), Inches(0.8), "489", 34, WHITE, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_TITLE)
txt(s, Inches(5.2), Inches(4.85), Inches(2.05), Inches(0.4), "Down ↓", 13, WHITE,
    align=PP_ALIGN.CENTER, font=F_BODY)
# top genler
txt(s, Inches(0.8), Inches(5.6), Inches(6.6), Inches(0.4), "Öne çıkan genler:", 16, NAVY, bold=True, font=F_BODY)
bullets(s, Inches(0.8), Inches(6.05), Inches(6.6), Inches(2.4),
    [
     "GPC3 (↑, logFC +3.5): bilinen HCC biyobelirteci",
     "TOP2A, AKR1B10 (↑): proliferasyon / metabolizma",
     "HAMP, CYP1A2, FCN3 (↓): karaciğer fonksiyon kaybı",
    ], size=16, gap=7)
# sağ: volcano
pic_fit(s, os.path.join(FIG, "volcano_plot.png"), Inches(7.8), Inches(2.0), Inches(7.5), Inches(6.2))
footer(s, 7)


# ==============================================================================
# SLAYT 8 — DGE GÖRSELLER
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "DGE: Top-5 Gen Dağılımları", kicker="GÖRSELLEŞTİRME")
pic_fit(s, os.path.join(FIG, "boxplot_top5_with_stats.png"), Inches(0.6), Inches(1.9), Inches(14.8), Inches(3.0))
pic_fit(s, os.path.join(FIG, "violin_top5.png"), Inches(0.6), Inches(5.0), Inches(14.8), Inches(3.0))
txt(s, Inches(0.6), Inches(8.05), Inches(14.8), Inches(0.5),
    "Kutu ve keman grafikleri Wilcoxon testi ile; tüm Top-5 gen iki grup arasında anlamlı fark gösterir (p < 0.0001).",
    14, GREY, italic=True, align=PP_ALIGN.CENTER, font=F_BODY)
footer(s, 8)


# ==============================================================================
# SLAYT 9 — SAĞKALIM
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "Sağkalım Analizi: MT1G", kicker="KAPLAN-MEIER & COX")
pic_fit(s, os.path.join(FIG, "km_MT1G.png"), Inches(0.7), Inches(2.0), Inches(7.4), Inches(5.6))
# sağ yorum
bullets(s, Inches(8.4), Inches(2.1), Inches(7.0), Inches(4),
    [
     "Sadece tümör örnekleri (n=115), OS (genel sağkalım) üzerinden.",
     "Top-5 up + Top-5 down = 10 gen test edildi.",
     "Anlamlı: MT1G (p=0.006), MT1H (p=0.023), AKR1B10 (p=0.039)",
     ("Yüksek MT1G → kötü sağkalım (Cox HR=3.35)", 1),
     ("MT1G/MT1H: metallotiyonein ailesi", 1),
     "Cox modeli: gen grubu + yaş + cinsiyet kovaryatları.",
     "RFS'de (düzeltilmiş eşleştirme) anlamlı gen yok.",
    ], size=18, gap=9)
rrect(s, Inches(8.4), Inches(6.75), Inches(7.0), Inches(1.1), GREEN)
txt(s, Inches(8.6), Inches(6.9), Inches(6.6), Inches(0.8),
    "MT1G: tümörde baskılanır; tümör-içi yüksek ifade kötü prognozla ilişkili",
    16, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, font=F_BODY)
footer(s, 9)


# ==============================================================================
# SLAYT 10 — WGCNA
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "WGCNA Ağ Analizi", kicker="CO-EXPRESSION")
pic_fit(s, os.path.join(FIG, "module_trait_heatmap.png"), Inches(0.7), Inches(2.0), Inches(6.2), Inches(5.2))
pic_fit(s, os.path.join(FIG, "network_brown_module.png"), Inches(7.0), Inches(2.0), Inches(5.0), Inches(5.2))
# sağ hub gen kutusu
rrect(s, Inches(12.2), Inches(2.0), Inches(3.2), Inches(5.2), LGREY)
txt(s, Inches(12.35), Inches(2.1), Inches(2.9), Inches(0.5), "Brown Modül Hub Genleri", 15, NAVY, bold=True, font=F_BODY)
hubs = ["GLYATL1", "TTC36", "PEX11G", "AADAT", "NAT2"]
yy = Inches(2.75)
for hname in hubs:
    rrect(s, Inches(12.35), yy, Inches(2.9), Inches(0.55), RED)
    txt(s, Inches(12.35), yy, Inches(2.9), Inches(0.55), hname, 15, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_BODY)
    yy += Inches(0.68)
txt(s, Inches(12.35), yy+Inches(0.05), Inches(2.9), Inches(1.2),
    "Karaciğer metabolizması & detoksifikasyon genleri", 13, GREY, italic=True, font=F_BODY)
txt(s, Inches(0.7), Inches(7.45), Inches(11), Inches(0.9),
    "6 co-expression modülü · β=5 · Brown modül (1636 gen) tümörle en güçlü korelasyonlu · "
    "her modülde 5 hub gen (kME + GS).",
    15, GREY, italic=True, font=F_BODY)
footer(s, 10)


# ==============================================================================
# SLAYT 11 — GO ENRICHMENT
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "GO Zenginleştirme Analizi", kicker="clusterProfiler")
pic_fit(s, os.path.join(FIG, "go_dotplot.png"), Inches(0.7), Inches(2.0), Inches(8.0), Inches(5.6))
bullets(s, Inches(9.0), Inches(2.1), Inches(6.4), Inches(4.5),
    [
     "Brown modül (1636 gen) için GO analizi.",
     "752 anlamlı Biyolojik Süreç (BP) terimi.",
     "En anlamlı süreçler:",
     ("Küçük molekül katabolizması", 1),
     ("Amino asit metabolizması", 1),
     ("Yağ asidi metabolizması", 1),
     ("Organik asit katabolizması", 1),
    ], size=18, gap=9)
rrect(s, Inches(9.0), Inches(6.4), Inches(6.4), Inches(1.5), NAVY)
txt(s, Inches(9.2), Inches(6.55), Inches(6.0), Inches(1.2),
    "Bu terimler karaciğerin temel metabolik fonksiyonları. Tümörle negatif korelasyon → "
    "hepatosit diferansiyasyon kaybının ağ düzeyinde kanıtı.",
    15, WHITE, anchor=MSO_ANCHOR.MIDDLE, font=F_BODY, line_sp=1.05)
footer(s, 11)


# ==============================================================================
# SLAYT 12 — ML SETUP'LAR
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "Makine Öğrenmesi: Özellik Setup'ları", kicker="FEATURE SELECTION")
setups = [
    ("Setup-1 · Biyolojik", "742 gen\n\nDGE up/down +\nWGCNA hub genleri\n\nBiyolojik anlam", BLUE),
    ("Setup-2 · İstatistiksel", "446 gen\n\nVaryans filtre +\nkorelasyon filtre\n(r < 0.9)", GREEN),
    ("Setup-3 · PCA", "35 bileşen\n\nPCA boyut indirgeme\n(%80 varyans)\n\nFeature extraction", ORANGE),
]
xx = Inches(0.9)
for h, b, c in setups:
    card(s, xx, Inches(2.1), Inches(4.5), Inches(3.4), h, b, c, body_size=16, head_size=18)
    xx += Inches(4.85)
# modeller
txt(s, Inches(0.9), Inches(5.9), Inches(14), Inches(0.5), "Üç eğiticili model · 5-Katlı Çapraz Doğrulama:", 18, NAVY, bold=True, font=F_BODY)
models = [("SVM", "Radial kernel\nMargin tabanlı"), ("Random Forest", "Ensemble bagging\nKarar ağaçları"), ("XGBoost", "Gradient boosting\nSıralı ağaçlar")]
xx = Inches(0.9)
for h, b in models:
    rrect(s, xx, Inches(6.5), Inches(4.5), Inches(1.6), LGREY)
    txt(s, xx+Inches(0.2), Inches(6.65), Inches(4.1), Inches(0.5), h, 18, NAVY, bold=True, font=F_BODY)
    txt(s, xx+Inches(0.2), Inches(7.15), Inches(4.1), Inches(0.9), b, 15, DARK, font=F_BODY)
    xx += Inches(4.85)
footer(s, 12)


# ==============================================================================
# SLAYT 13 — ML SONUÇLAR
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "ML Sonuçları: 3 Model × 3 Setup", kicker="DEĞERLENDİRME")
# tablo
rows = [
    ("Setup", "Model", "Acc", "F1", "AUC"),
    ("Biyolojik", "SVM", "0.946", "0.961", "0.975"),
    ("Biyolojik", "Random Forest", "0.958", "0.969", "0.980"),
    ("Biyolojik", "XGBoost", "0.952", "0.965", "0.973"),
    ("İstatistiksel", "Random Forest", "0.958", "0.970", "0.981"),
    ("PCA", "SVM", "0.958", "0.969", "0.967"),
]
tx = Inches(0.8); ty = Inches(2.0)
col_w = [Inches(2.4), Inches(2.7), Inches(1.3), Inches(1.3), Inches(1.3)]
rh = Inches(0.62)
for ri, row in enumerate(rows):
    cx = tx
    head = (ri == 0)
    bg = NAVY if head else (LGREY if ri % 2 else WHITE)
    for ci, cell in enumerate(row):
        rect(s, cx, ty, col_w[ci], rh, bg)
        txt(s, cx+Inches(0.1), ty, col_w[ci]-Inches(0.2), rh, cell,
            15, WHITE if head else DARK, bold=head,
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER if ci >= 2 else PP_ALIGN.LEFT, font=F_BODY)
        cx += col_w[ci]
    ty += rh
txt(s, Inches(0.8), ty+Inches(0.1), Inches(8.2), Inches(1.0),
    "Hasta-bazlı CV + sızıntısız PCA. Tüm 9 kombinasyon: Acc > %93, AUC > 0.95.\nEn iyi AUC: İstatistiksel + Random Forest (0.981).",
    15, GREY, italic=True, font=F_BODY, line_sp=1.1)
# ROC
pic_fit(s, os.path.join(FIG, "roc_curves_Setup1_Biological.png"), Inches(9.4), Inches(2.0), Inches(6.0), Inches(5.0))
txt(s, Inches(9.4), Inches(7.1), Inches(6.0), Inches(0.5),
    "Setup-1 (Biyolojik) ROC eğrileri", 14, GREY, italic=True, align=PP_ALIGN.CENTER, font=F_BODY)
footer(s, 13)


# ==============================================================================
# SLAYT 14 — SONUÇLAR
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "Sonuçlar ve Değerlendirme", kicker="ÖZET")
res = [
    ("705 DEG", "216 up / 489 down. Down baskınlığı → hepatosit diferansiyasyon kaybı.", RED),
    ("MT1G", "OS'de anlamlı prognostik gen (p=0.006, HR=3.35). GPC3 ise en güçlü tanısal DEG.", GREEN),
    ("WGCNA", "6 modül; brown modül karaciğer metabolizmasını temsil eder, tümörle en korelasyonlu.", ORANGE),
    ("GO", "752 BP terimi; amino asit, yağ asidi, küçük molekül metabolizması.", LBLUE),
    ("ML", "3 model × 3 setup, tümü AUC > 0.95. Biyolojik ≈ istatistiksel feature selection.", BLUE),
]
yy = Inches(1.95)
for h, b, c in res:
    rrect(s, Inches(0.9), yy, Inches(2.7), Inches(0.92), c)
    txt(s, Inches(0.9), yy, Inches(2.7), Inches(0.92), h, 19, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_TITLE)
    txt(s, Inches(3.8), yy, Inches(11.4), Inches(0.92), b, 17, DARK,
        anchor=MSO_ANCHOR.MIDDLE, font=F_BODY)
    yy += Inches(1.06)
rrect(s, Inches(0.9), yy+Inches(0.02), Inches(14.3), Inches(0.8), NAVY)
txt(s, Inches(1.1), yy+Inches(0.02), Inches(13.9), Inches(0.8),
    "Genel sonuç: HCC tümörü normal karaciğer fonksiyonlarını kaybeder, onkojenik yolakları aktive eder; "
    "gen ifadesi tümör/normal ayrımında güçlü bir sinyaldir.",
    16, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, font=F_BODY, line_sp=1.05)
footer(s, 14)


# ==============================================================================
# SLAYT 15 — TEŞEKKÜR
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, NAVY)
rect(s, Inches(0), Inches(4.3), SW, Pt(4), RED)
txt(s, Inches(1), Inches(3.2), Inches(14), Inches(1.3),
    "Dinlediğiniz için teşekkürler!", 48, WHITE, bold=True, font=F_TITLE, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.6), Inches(14), Inches(0.6),
    "Mert GÜLER · BLM3810 Biyoenformatiğe Giriş · GSE76427 HCC Projesi",
    18, LBLUE, font=F_BODY, align=PP_ALIGN.CENTER)


# ==============================================================================
# SLAYT 16 — KAYNAKÇA
# ==============================================================================
s = add_slide()
rect(s, Inches(0), Inches(0), SW, SH, WHITE)
content_header(s, "Kaynakça")
refs = [
    ("Veri Seti", "GSE76427 — NCBI Gene Expression Omnibus (GEO)"),
    ("limma", "Ritchie et al. (2015), Nucleic Acids Research 43(7):e47"),
    ("WGCNA", "Langfelder & Horvath (2008), BMC Bioinformatics 9:559"),
    ("clusterProfiler", "Yu et al. (2012), OMICS 16(5):284-287"),
    ("survival", "Therneau & Grambsch (2000), Modeling Survival Data, Springer"),
    ("caret / xgboost", "Kuhn (2008) JSS 28(5); Chen & Guestrin (2016) KDD"),
    ("GPC3", "Capurro et al. (2003), Gastroenterology 125(1):89-97"),
]
yy = Inches(2.2)
for k, v in refs:
    rrect(s, Inches(1.0), yy, Inches(3.4), Inches(0.7), NAVY)
    txt(s, Inches(1.0), yy, Inches(3.4), Inches(0.7), k, 16, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_BODY)
    txt(s, Inches(4.7), yy, Inches(10.5), Inches(0.7), v, 16, DARK,
        anchor=MSO_ANCHOR.MIDDLE, font=F_BODY)
    yy += Inches(0.82)
txt(s, Inches(1.0), yy+Inches(0.1), Inches(14), Inches(0.5),
    "Analizler R 4.3.3 ortamında gerçekleştirilmiştir. Tüm kod ve görseller proje deposunda mevcuttur.",
    14, GREY, italic=True, font=F_BODY)
footer(s, 16)


# ==============================================================================
prs.save(OUT)
print("Sunum kaydedildi:", OUT)
print("Toplam slayt:", len(prs.slides._sldIdLst))
